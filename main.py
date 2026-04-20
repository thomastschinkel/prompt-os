import traceback
import tkinter as tk
import subprocess
import threading
import sys
import json
import os
import re
import tkinter.font as tkfont
import customtkinter as ctk

stop_event = threading.Event()
ai_start_index = "1.0"
chat_history = []

# Global configuration paths (loaded on demand to speed up startup)
_SETTINGS_PATH = None
_CHATS_DIR = None

def get_settings_path():
    global _SETTINGS_PATH
    if _SETTINGS_PATH is None:
        from src.utils import get_config_path
        _SETTINGS_PATH = get_config_path("settings.json")
    return _SETTINGS_PATH

def get_chats_dir():
    global _CHATS_DIR
    if _CHATS_DIR is None:
        from src.utils import get_config_path
        path = get_config_path("settings.json")
        _CHATS_DIR = os.path.join(os.path.dirname(path), "chats")
        os.makedirs(_CHATS_DIR, exist_ok=True)
    return _CHATS_DIR

current_chat_file = None


def load_settings():
    try:
        path = get_settings_path()
        if os.path.exists(path):
            with open(path, "r", encoding="utf-8") as f:
                settings = json.load(f)
                if "enabled_tools" not in settings:
                    settings["enabled_tools"] = ["BASH", "POWERSHELL", "FILE_HANDLER", "EXEC_PY", "SEARCH_WEB", "READ_FILE",
                                                 "USE_BROWSER", "CLIPBOARD_MANAGER", "INTERPRET_SCREENSHOT", "GLOB", "GREP"]
                return settings
    except Exception:
        pass
    return {
        "provider": "GitHubAI",
        "model": "openai/gpt-4o-mini",
        "mode": "FAST",
        "ollama_base_url": "http://localhost:11434/v1",
        "lmstudio_base_url": "http://localhost:1234/v1",
        "enabled_tools": ["BASH", "POWERSHELL", "FILE_HANDLER", "EXEC_PY", "SEARCH_WEB", "READ_FILE", "USE_BROWSER",
                          "CLIPBOARD_MANAGER", "INTERPRET_SCREENSHOT", "GLOB", "GREP"]
    }


def save_settings(provider, model, mode, enabled_tools=None):
    try:
        current = load_settings()
        settings = dict(current)
        settings["provider"] = provider
        settings["model"] = model
        settings["mode"] = mode
        settings["enabled_tools"] = enabled_tools if enabled_tools is not None else current.get("enabled_tools", [])
        settings.setdefault("ollama_base_url", "http://localhost:11434/v1")
        settings.setdefault("lmstudio_base_url", "http://localhost:1234/v1")
        with open(get_settings_path(), "w", encoding="utf-8") as f:
            json.dump(settings, f, indent=4)
    except Exception:
        pass


def update_status(text):
    root.after(0, lambda: status_label.configure(text=text))


def insert_inline_markdown(text_widget, line):
    pattern = r"(`[^`]+`|\*\*[^*]+\*\*|\*[^*]+\*)"
    parts = re.split(pattern, line)

    for part in parts:
        if not part:
            continue
        if part.startswith("`") and part.endswith("`") and len(part) >= 2:
            text_widget.insert(tk.END, part[1:-1], ("md_inline_code",))
        elif part.startswith("**") and part.endswith("**") and len(part) >= 4:
            text_widget.insert(tk.END, part[2:-2], ("md_bold",))
        elif part.startswith("*") and part.endswith("*") and len(part) >= 2:
            text_widget.insert(tk.END, part[1:-1], ("md_italic",))
        else:
            text_widget.insert(tk.END, part)


def render_markdown_to_textbox(text_widget, text):
    in_code_block = False
    for raw_line in text.splitlines():
        stripped = raw_line.strip()

        if stripped.startswith("```"):
            in_code_block = not in_code_block
            continue

        if in_code_block:
            text_widget.insert(tk.END, f"{raw_line}\n", ("md_code_block",))
            continue

        if stripped.startswith("### "):
            text_widget.insert(tk.END, stripped[4:] + "\n", ("md_h3",))
            continue
        if stripped.startswith("## "):
            text_widget.insert(tk.END, stripped[3:] + "\n", ("md_h2",))
            continue
        if stripped.startswith("# "):
            text_widget.insert(tk.END, stripped[2:] + "\n", ("md_h1",))
            continue

        if re.match(r"^[-*+]\s+", stripped):
            text_widget.insert(tk.END, "• ", ("md_list",))
            insert_inline_markdown(text_widget, re.sub(r"^[-*+]\s+", "", stripped))
            text_widget.insert(tk.END, "\n")
            continue

        if re.match(r"^\d+\.\s+", stripped):
            insert_inline_markdown(text_widget, stripped)
            text_widget.insert(tk.END, "\n")
            continue

        if stripped.startswith("> "):
            text_widget.insert(tk.END, "| ", ("md_quote",))
            insert_inline_markdown(text_widget, stripped[2:])
            text_widget.insert(tk.END, "\n")
            continue

        insert_inline_markdown(text_widget, raw_line)
        text_widget.insert(tk.END, "\n")


def configure_markdown_tags(text_widget):
    internal_text = text_widget._textbox
    base_font = tkfont.Font(font=internal_text.cget("font"))
    base_family = base_font.actual("family")
    base_size = base_font.actual("size")
    size_offset = -1 if base_size < 0 else 1

    text_widget._md_fonts = {
        "h1": tkfont.Font(family=base_family, size=base_size + (4 * size_offset), weight="bold"),
        "h2": tkfont.Font(family=base_family, size=base_size + (2 * size_offset), weight="bold"),
        "h3": tkfont.Font(family=base_family, size=base_size + (1 * size_offset), weight="bold"),
        "bold": tkfont.Font(family=base_family, size=base_size, weight="bold"),
        "italic": tkfont.Font(family=base_family, size=base_size, slant="italic"),
        "inline_code": tkfont.Font(family="Consolas", size=base_size),
        "code_block": tkfont.Font(family="Consolas", size=base_size - (1 * size_offset))
    }

    internal_text.tag_config("md_h1", foreground="#ffffff", spacing1=10, spacing3=8, font=text_widget._md_fonts["h1"])
    internal_text.tag_config("md_h2", foreground="#ffffff", spacing1=8, spacing3=6, font=text_widget._md_fonts["h2"])
    internal_text.tag_config("md_h3", foreground="#ffffff", spacing1=6, spacing3=4, font=text_widget._md_fonts["h3"])
    internal_text.tag_config("md_bold", foreground="#ffffff", font=text_widget._md_fonts["bold"])
    internal_text.tag_config("md_italic", foreground="#e0e0e0", font=text_widget._md_fonts["italic"])
    internal_text.tag_config("md_inline_code", background="#2a2a3a", foreground="#f6d365", font=text_widget._md_fonts["inline_code"])
    internal_text.tag_config("md_code_block", background="#1a1a28", foreground="#c9f0ff", lmargin1=10, lmargin2=10, font=text_widget._md_fonts["code_block"])
    internal_text.tag_config("md_list", foreground="#e0e0e0")
    internal_text.tag_config("md_quote", foreground="#a9a9c5", lmargin1=10, lmargin2=10)
    internal_text.tag_config("md_error", foreground="#ff5555")
    internal_text.tag_config("md_thinking", foreground="#888888", font=text_widget._md_fonts["italic"])
    internal_text.tag_config("user_label", foreground="#7fb5ff", spacing1=15, spacing3=5, font=text_widget._md_fonts["bold"])
    internal_text.tag_config("ai_label", foreground="#00e676", spacing1=15, spacing3=5, font=text_widget._md_fonts["h3"])
    internal_text.tag_config("user_msg", foreground="#c0c0d0", lmargin1=10, lmargin2=10)


def get_image_description(img_str, provider, model, keys):
    from openai import OpenAI
    from src.utils import get_config_path
    client = None
    fallback_model = model
    settings = load_settings()

    if provider == "Groq":
        client = OpenAI(base_url="https://api.groq.com/openai/v1", api_key=keys.get("GROQ_API_KEY", ""))
        fallback_model = "meta-llama/llama-4-scout-17b-16e-instruct"
    elif provider == "GitHubAI":
        client = OpenAI(base_url="https://models.github.ai/inference", api_key=keys.get("GITHUB_TOKEN", ""))
        fallback_model = "Phi-4-multimodal-instruct"
    elif provider == "OpenRoute":
        client = OpenAI(base_url="https://openrouter.ai/api/v1", api_key=keys.get("OPENROUTE_API_KEY", ""))
        fallback_model = "gemma-4-31b-it"
    elif provider == "Unclose":
        client = OpenAI(base_url="https://qwen-vl.ai.unturf.com/v1/", api_key="is_free")
        fallback_model = model
    elif provider == "OpenAI":
        client = OpenAI(api_key=keys.get("OPENAI_API_KEY", ""))
        fallback_model = "gpt-4.1"
    elif provider == "Anthropic":
        client = OpenAI(base_url="https://api.anthropic.com/v1/", api_key=keys.get("ANTHROPIC_API_KEY", ""))
        fallback_model = "claude-haiku-4-5"
    elif provider == "Google":
        client = OpenAI(base_url="https://generativelanguage.googleapis.com/v1beta/openai/",
                        api_key=keys.get("GOOGLE_API_KEY", ""))
        fallback_model = model
    elif provider == "Ollama":
        client = OpenAI(base_url=settings.get("ollama_base_url", "http://localhost:11434/v1"),
                        api_key=keys.get("OLLAMA_API_KEY", "ollama"))
        fallback_model = model
    elif provider == "LMStudio":
        client = OpenAI(base_url=settings.get("lmstudio_base_url", "http://localhost:1234/v1"),
                        api_key=keys.get("LMSTUDIO_API_KEY", "lm-studio"))
        fallback_model = model

    if not client:
        return "Error: Unknown provider for image interpretation."

    messages = [{
        "role": "user",
        "content": [
            {"type": "text",
             "text": "Describe this image in extreme detail. Identify objects, colors, text, and overall composition."},
            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_str}"}}
        ]
    }]
    try:
        res = client.chat.completions.create(model=model, messages=messages)
        return res.choices[0].message.content
    except Exception:
        try:
            res = client.chat.completions.create(model=fallback_model, messages=messages)
            return res.choices[0].message.content
        except Exception as e:
            return f"Error describing image (both models failed): {e}"


def update_answer_box(text, color_tag=None, done_event=None):
    def _update():
        ai_answer_box.configure(state="normal")
        try:
            ai_answer_box._textbox.delete("ai_start", tk.END)
        except tk.TclError:
            pass
        if color_tag:
            t = text if text.endswith("\n") else text + "\n"
            tags = (color_tag, "left_justify")
            ai_answer_box.insert(tk.END, t, tags)
        else:
            ai_answer_box._textbox.insert(tk.END, "\n", ("left_justify",))
            render_markdown_to_textbox(ai_answer_box, text)
            ai_answer_box._textbox.tag_add("left_justify", "ai_start", tk.END)
        ai_answer_box.configure(state="disabled")
        ai_answer_box.see(tk.END)
        if done_event:
            done_event.set()

    root.after(0, _update)


class ExpandableTool(ctk.CTkFrame):
    def __init__(self, master, title, details, **kwargs):
        import customtkinter as ctk
        super().__init__(master, fg_color="#1e1e2e", **kwargs)
        self.expanded = False
        self.details = details
        
        self.header_btn = ctk.CTkButton(self, text=f"  ▶ {title}", anchor="w", fg_color="transparent", 
                                        hover_color="#2a2a3a", text_color="#7a7a8c", 
                                        font=ctk.CTkFont(size=13), height=28, command=self.toggle)
        self.header_btn.pack(fill="x", pady=0)
        
        self.content_lbl = ctk.CTkTextbox(self, fg_color="#14141e", text_color="#a9a9c5", 
                                          font=ctk.CTkFont(family="Consolas", size=12),
                                          height=80, wrap="word", border_width=1, border_color="#2b2d35", corner_radius=8)
        self.content_lbl.insert("1.0", self.details)
        self.content_lbl.configure(state="disabled")
        
    def destroy(self):
        # Manually destroy the textbox first to prevent TclError in its after loop
        try:
            if hasattr(self, "content_lbl"):
                # Cancel any pending after callbacks in CustomTkinter Textbox if possible
                self.content_lbl._textbox.after_cancel(self.content_lbl._textbox)
        except:
            pass
        super().destroy()
        
    def toggle(self):
        if not self.winfo_exists(): return
        self.expanded = not self.expanded
        current_text = self.header_btn.cget("text")
        if self.expanded:
            self.header_btn.configure(text=current_text.replace("▶", "▼", 1))
            self.content_lbl.pack(fill="x", padx=(30, 20), pady=(0, 6))
        else:
            self.header_btn.configure(text=current_text.replace("▼", "▶", 1))
            self.content_lbl.pack_forget()

def insert_tool_widget(title, details, done_event=None):
    def _insert():
        ai_answer_box.configure(state="normal")
        try:
            ai_answer_box._textbox.delete("ai_start", tk.END)
        except tk.TclError:
            pass
        
        tool_frame = ExpandableTool(ai_answer_box, title, details)
        
        if ai_answer_box.get("end-2c", "end-1c") != "\n":
            ai_answer_box.insert(tk.END, "\n", ("left_justify",))
        ai_answer_box._textbox.window_create(tk.END, window=tool_frame)
        ai_answer_box.insert(tk.END, "\n", ("left_justify",))
        
        ai_answer_box.insert(tk.END, "\u200b", ("left_justify",))
        ai_answer_box._textbox.mark_set("ai_start", "end-1c")
        ai_answer_box._textbox.mark_gravity("ai_start", "left")
        ai_answer_box.configure(state="disabled")
        ai_answer_box.see(tk.END)
        if done_event:
            done_event.set()
    root.after(0, _insert)

def toggle_send_stop(generating=True):
    try:
        if generating:
            send_button.pack_forget()
            stop_button.pack(side="right", padx=(10, 0), anchor="s")
        else:
            stop_button.pack_forget()
            send_button.pack(side="right", padx=(10, 0), anchor="s")
    except Exception:
        pass

def stop_task():
    try:
        if stop_button.winfo_ismapped():
            stop_event.set()
            update_status("Stopping...")
            toggle_send_stop(False)
    except Exception:
        pass

is_thinking = False
def start_thinking_animation():
    global is_thinking
    if is_thinking: return
    is_thinking = True
    update_answer_box("Thinking...\n", "md_thinking")

def stop_thinking_animation():
    global is_thinking
    is_thinking = False

def render_chat_history():
    ai_answer_box.configure(state="normal")
    
    # Explicitly destroy children to prevent TclError from pending callbacks in CustomTkinter
    for window_name in ai_answer_box._textbox.window_names():
        try:
            w = ai_answer_box._textbox.nametowidget(window_name)
            w.destroy()
        except Exception:
            pass
            
    ai_answer_box.delete("1.0", tk.END)
    ai_answer_box._textbox.tag_configure("right_justify", justify="right")
    ai_answer_box._textbox.tag_configure("left_justify", justify="left")
    for msg in chat_history:
        if ai_answer_box.index("end-1c") != "1.0" and msg["role"] != "tool":
            ai_answer_box.insert(tk.END, "\n\n")
            
        if msg["role"] == "user":
            import customtkinter as ctk
            bubble_frame = ctk.CTkFrame(ai_answer_box, fg_color="#2b2d35", corner_radius=12, border_width=1, border_color="#3d3f4b")
            lbl = ctk.CTkLabel(bubble_frame, text=msg["content"].strip(), font=ctk.CTkFont(size=14, family="Helvetica"), text_color="#e0e0e0", justify="left", wraplength=450)
            lbl.pack(padx=12, pady=10)
            
            ai_answer_box._textbox.window_create(tk.END, window=bubble_frame)
            ai_answer_box.insert(tk.END, "\n", ("right_justify",))
            ai_answer_box.insert(tk.END, "\u200b", ("left_justify",))
        elif msg["role"] == "tool":
            tool_frame = ExpandableTool(ai_answer_box, msg["title"], msg["details"])
            if ai_answer_box.get("end-2c", "end-1c") != "\n":
                ai_answer_box.insert(tk.END, "\n", ("left_justify",))
            ai_answer_box._textbox.window_create(tk.END, window=tool_frame)
            ai_answer_box.insert(tk.END, "\n", ("left_justify",))
            ai_answer_box.insert(tk.END, "\u200b", ("left_justify",))
            ai_answer_box._textbox.mark_set("ai_start", "end-1c")
        else:
            render_markdown_to_textbox(ai_answer_box, msg["content"].strip())
    ai_answer_box.configure(state="disabled")
    ai_answer_box.see(tk.END)

current_llm = None

def clear_chat():
    global current_llm, current_chat_file
    current_llm = None
    current_chat_file = None
    chat_history.clear()
    ai_answer_box.configure(state="normal")
    ai_answer_box.delete("1.0", tk.END)
    ai_answer_box.configure(state="disabled")
    update_status("New Chat started")
    update_sidebar_list()

def save_current_chat():
    global current_chat_file
    if not chat_history: return
    
    if current_chat_file is None:
        first_msg = next((m["content"] for m in chat_history if m["role"] == "user"), "New Chat")
        safe_name = "".join([c for c in first_msg[:20] if c.isalnum() or c in " _-"]).strip()
        current_chat_file = os.path.join(get_chats_dir(), f"{safe_name}_{int(threading.current_thread().ident)}.json")

    with open(current_chat_file, "w", encoding="utf-8") as f:
        json.dump(chat_history, f, indent=2)

def load_chat_file(filepath):
    global current_chat_file, chat_history, current_llm
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            data = json.load(f)
            chat_history.clear()
            chat_history.extend(data)
            current_chat_file = filepath
            current_llm = None # Reset LLM to reload history correctly
            render_chat_history()
            update_status(f"Loaded: {os.path.basename(filepath)}")
    except Exception as e:
        update_status(f"Error loading chat: {e}")

def delete_chat(filepath):
    global current_chat_file
    try:
        if os.path.exists(filepath):
            os.remove(filepath)
            if current_chat_file == filepath:
                clear_chat()
            else:
                update_sidebar_list()
            update_status(f"Deleted: {os.path.basename(filepath)}")
    except Exception as e:
        update_status(f"Error deleting chat: {e}")

def update_sidebar_list():
    for child in history_scroll.winfo_children():
        child.destroy()
        
    chats_dir = get_chats_dir()
    files = sorted([f for f in os.listdir(chats_dir) if f.endswith(".json")], 
                   key=lambda x: os.path.getmtime(os.path.join(chats_dir, x)), reverse=True)
    
    for f in files:
        full_path = os.path.join(chats_dir, f)
        display_name = f.replace(".json", "")
        # Remove timestamp/id if present (heuristic)
        display_name = re.sub(r"_\d+$", "", display_name)
        
        chat_item = ctk.CTkFrame(history_scroll, fg_color="transparent")
        chat_item.pack(fill="x", pady=1, padx=5)
        
        # Pack delete button first to the right so it stays fixed
        del_btn = ctk.CTkButton(chat_item, text="✕", width=24, height=24, fg_color="transparent",
                                hover_color="#c42b2b", text_color="#7a7a8c",
                                font=ctk.CTkFont(size=12, weight="bold"),
                                command=lambda p=full_path: delete_chat(p))
        del_btn.pack(side="right", padx=(2, 0))
        
        btn = ctk.CTkButton(chat_item, text=display_name, anchor="w", fg_color="transparent", 
                            hover_color="#2a2a3a", text_color="#a9a9c5", height=28, width=0,
                            command=lambda p=full_path: load_chat_file(p))
        btn.pack(side="left", fill="x", expand=True)

def copy_last():
    for msg in reversed(chat_history):
        if msg["role"] == "ai":
            root.clipboard_clear()
            root.clipboard_append(msg["content"])
            update_status("Copied to clipboard")
            return
    update_status("Nothing to copy")

def regenerate():
    try:
        if not send_button.winfo_ismapped() or not chat_history:
            return
        if chat_history[-1]["role"] == "ai":
            chat_history.pop()
        if chat_history and chat_history[-1]["role"] == "user":
            global current_llm
            if current_llm and current_llm.history:
                 while current_llm.history and current_llm.history[-1]["role"] != "user":
                     current_llm.history.pop()
                 if current_llm.history and not current_llm.history[-1].get("content", "").startswith("Command output:"):
                     current_llm.history.pop()
            threading.Thread(target=handle_task, args=(True,), daemon=True).start()
    except Exception:
        pass

def thread_handle_task():
    try:
        if not send_button.winfo_ismapped(): return
    except Exception:
        pass
    threading.Thread(target=handle_task, daemon=True).start()

def resize_textbox(event=None):
    if not isinstance(text_input, ctk.CTkTextbox): return
    text = text_input.get("1.0", "end-1c")
    lines = text.count("\n") + 1
    new_height = min(max(1, lines), 6) * 20 + 20
    text_input.configure(height=new_height)

def handle_task(is_regenerate=False):
    global stop_event, ai_start_index

    if not is_regenerate:
        user_input = text_input.get("1.0", tk.END).strip() if isinstance(text_input, ctk.CTkTextbox) else text_input.get().strip()
        if not user_input.strip():
            return
        if isinstance(text_input, ctk.CTkTextbox):
            text_input.delete("1.0", tk.END)
            resize_textbox()
        else:
            text_input.delete(0, tk.END)
        chat_history.append({"role": "user", "content": user_input})
    else:
        user_input = chat_history[-1]["content"] if chat_history else ""
        if not user_input: return

    from src.utils import get_config_path, decode_output
    setup_event = threading.Event()
    def _setup_ui():
        render_chat_history()
        ai_answer_box.configure(state="normal")
        if ai_answer_box.index("end-1c") != "1.0" and ai_answer_box.get("end-2c", "end-1c") != "\n":
            ai_answer_box.insert(tk.END, "\n", ("left_justify",))
        ai_answer_box.insert(tk.END, "\u200b", ("left_justify",)) # Zero-width char as anchor
        ai_answer_box._textbox.mark_set("ai_start", "end-2c")
        ai_answer_box._textbox.mark_gravity("ai_start", "left")
        ai_answer_box.configure(state="disabled")
        toggle_send_stop(True)
        stop_button.configure(fg_color="#c42b2b", hover_color="#a82525")
        setup_event.set()

    root.after(0, _setup_ui)
    setup_event.wait()

    current_stop_event = threading.Event()
    stop_event = current_stop_event
    my_stop_event = current_stop_event

    global current_llm
    from src.ai import LLM
    provider = provider_var.get()
    model = model_var.get()
    mode = mode_var.get()

    if current_llm is None or current_llm.provider != provider or current_llm.model_name != model:
        current_llm = LLM(provider=provider, model_name=model, mode=mode, stop_event=current_stop_event)
    else:
        current_llm.stop_event = current_stop_event
        current_llm.mode = mode

    llm = current_llm
    root.after(0, start_thinking_animation)
    settings = load_settings()
    enabled_tools = settings.get("enabled_tools", [])
    disabled_tools = [t for t in
                      ["BASH", "POWERSHELL", "FILE_HANDLER", "EXEC_PY", "SEARCH_WEB", "READ_FILE", "USE_BROWSER", "CLIPBOARD_MANAGER",
                       "INTERPRET_SCREENSHOT", "GLOB", "GREP"] if t not in enabled_tools]

    aborted = False

    def on_finish():
        root.after(0, toggle_send_stop, False)
        root.after(0, lambda: stop_button.configure(fg_color="#3a3a55", hover_color="#4a4a65"))

    prompt_input = user_input
    if not is_regenerate and len(chat_history) <= 1:
        prompt_input += "\n\n[SYSTEM: This is the very first turn of a new chat. You MUST provide a 'title' field in your JSON response.]"

    response = llm.generate_response(prompt_input, status_callback=update_status)

    # Chat Title Logic
    global current_chat_file
    generated_title = response.get("title")
    if not is_regenerate and len(chat_history) <= 2 and current_chat_file is None:
        try:
            # If no title yet, we use a temporary "New Chat" name but will try to update it later
            title_text = generated_title or "New Chat"
            title_text = title_text.strip().strip('"').strip("'")
            safe_title = "".join([c for c in title_text if c.isalnum() or c in " _-"]).strip()
            if not safe_title: safe_title = "Chat"
            current_chat_file = os.path.join(get_chats_dir(), f"{safe_title}_{int(threading.current_thread().ident)}.json")
            save_current_chat()
            root.after(0, update_sidebar_list)
        except Exception:
            pass

    if current_stop_event.is_set():
        on_finish()
        return

    loop_count = 0
    while response.get("status", "y") != "y":
        if current_stop_event.is_set():
            break

        root.after(0, stop_thinking_animation)
        loop_count += 1
        
        tool = response.get("tool", "")
        thought = response.get("response", "").strip() or response.get("thought", "").strip()
        
        if loop_count % 3 == 0 and thought:
            event1 = threading.Event()
            insert_tool_widget("Thinking...", thought, event1)
            chat_history.append({"role": "tool", "title": "Thinking...", "details": thought})
            event1.wait()
            
        if tool:
            event2 = threading.Event()
            # Standardize arguments (merge top-level and 'args' field)
            tool_args = response.get("args") if isinstance(response.get("args"), dict) else {}
            for k, v in response.items():
                if k not in ["response", "tool", "memory", "status", "thought", "args"]:
                    if k not in tool_args: tool_args[k] = v
            
            # Format display text (show key: value)
            display_parts = []
            # Priority keys to show first
            priority_keys = ["command", "query", "url", "path", "pattern", "code", "action", "input", "tool_input"]
            seen_keys = set()
            for pk in priority_keys:
                if pk in tool_args and tool_args[pk]:
                    display_parts.append(f"{pk}: {tool_args[pk]}")
                    seen_keys.add(pk)
            # Add any other parameters
            for k, v in tool_args.items():
                if k not in seen_keys and k != "thought" and v:
                    display_parts.append(f"{k}: {v}")
            
            details_text = "\n".join(display_parts) if display_parts else "No parameters provided"
                
            insert_tool_widget(f"Using Tool: {tool}", details_text, event2)
            chat_history.append({"role": "tool", "title": f"Using Tool: {tool}", "details": details_text})
            event2.wait()
            
            # Sync back to response object for existing execution blocks
            for k, v in tool_args.items():
                if k not in response: response[k] = v

        update_status(f"Running Tool: {tool}" if tool else "Thinking")
        output = ""

        if tool and tool != "ANSWER" and tool not in enabled_tools:
            output = f"Error: Tool '{tool}' is disabled by the administrator. DO NOT try to use this tool again in this conversation. If you have no other enabled tools to fulfill the request, inform the user and finish."
        # Tool execution check
        elif tool == "BASH":
            try:
                creation_flags = 0
                if os.name == 'nt':
                     creation_flags = subprocess.CREATE_NO_WINDOW
                cmd_input = response.get("input", response.get("tool_input", response.get("command", "")))
                result = subprocess.run(["bash", "-c", cmd_input], capture_output=True, text=False,
                                        timeout=None, creationflags=creation_flags)
                stdout = decode_output(result.stdout or b"")
                stderr = decode_output(result.stderr or b"")
                output = f"{stdout} | {stderr}".strip()
            except Exception as e:
                output = f"Error: {e}"

        elif tool == "POWERSHELL":
            try:
                creation_flags = 0
                if os.name == 'nt':
                     creation_flags = subprocess.CREATE_NO_WINDOW
                cmd_input = response.get("input", response.get("tool_input", response.get("command", "")))
                result = subprocess.run(["powershell.exe", "-Command", cmd_input], capture_output=True, text=False,
                                        timeout=None, creationflags=creation_flags)
                stdout = decode_output(result.stdout or b"")
                stderr = decode_output(result.stderr or b"")
                output = f"{stdout} | {stderr}".strip()
            except Exception as e:
                output = f"Error: {e}"

        elif tool == "FILE_HANDLER":
            path = response.get("path", response.get("file_path", response.get("input", response.get("tool_input", "")))).strip('"')
            mode = response.get("mode", "r")
            content_to_write = response.get("content", response.get("text", response.get("data", "")))
            if mode != "r":
                try:
                    with open(path, mode, encoding="utf-8") as f:
                        f.write(content_to_write)
                    output = f"File written to {path} with mode {mode}"
                except Exception:
                    output = f"Error writing to file {path} with mode {mode}"
            else:
                try:
                    with open(path, "r", encoding="utf-8") as f:
                        content = f.read()
                    output = f"Content of {path}: {content} with mode {mode}"
                except Exception as e:
                    output = f"Error reading file {path}: {str(e)}"

        elif tool == "EXEC_PY":
            try:
                from io import StringIO
                py_code = response.get("code", response.get("input", response.get("tool_input", "")))
                new_stdout, new_stderr = StringIO(), StringIO()
                old_stdout, old_stderr = sys.stdout, sys.stderr
                sys.stdout, sys.stderr = new_stdout, new_stderr
                try:
                    exec(py_code, {"__builtins__": __builtins__}, {})
                    out = f"STDOUT: {new_stdout.getvalue()}\nSTDERR: {new_stderr.getvalue()}"
                finally:
                    sys.stdout, sys.stderr = old_stdout, old_stderr
                output = out.strip() or "Executed (no output)"
            except Exception as e:
                output = f"ERROR: {str(e)}"

        elif tool == "SEARCH_WEB":
            search_query = response.get("query", response.get("input", response.get("tool_input", "")))
            if not search_query:
                output = "Error: Search query is empty. Please provide a non-empty 'query' or 'input'."
            else:
                from src.utils import search_web
                output = search_web(search_query, max_results=response.get("max_results", 3))

        elif tool == "READ_FILE":
            path = response.get("path", response.get("file_path", response.get("input", response.get("tool_input", "")))).strip('"')
            if not os.path.exists(path):
                output = f"Error: File not found at {path}"
            else:
                ext = path.lower().split('.')[-1]
                try:
                    if ext == "pdf":
                        import fitz
                        with fitz.open(path) as doc:
                            output = "\n".join(page.get_text() for page in doc)
                    elif ext == "docx":
                        import docx
                        output = "\n".join(p.text for p in docx.Document(path).paragraphs)
                    elif ext in ["png", "jpg", "jpeg", "webp", "bmp"]:
                        import base64
                        with open(path, "rb") as f:
                            img_str = base64.b64encode(f.read()).decode("utf-8")
                        try:
                            with open(get_config_path("keys.json"), "r", encoding="utf-8") as f:
                                keys = json.load(f)
                        except:
                            keys = {}
                        output = get_image_description(img_str, provider, model, keys)
                    elif ext == "pptx":
                        import pptx
                        prs = pptx.Presentation(path)
                        txt = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"): txt.append(shape.text)
                        output = "\n".join(txt)
                    elif ext in ["xlsx", "xls"]:
                        import pandas as pd
                        output = pd.read_excel(path).to_markdown()
                    elif ext == "csv":
                        import pandas as pd
                        output = pd.read_csv(path).to_markdown()
                    elif ext in ["html", "htm"]:
                        from bs4 import BeautifulSoup
                        with open(path, "r", encoding="utf-8") as f:
                            soup = BeautifulSoup(f.read(), "lxml")
                        for s in soup(["script", "style"]): s.decompose()
                        output = "\n".join(l.strip() for l in soup.get_text(separator="\n").splitlines() if l.strip())
                    else:
                        with open(path, "r", encoding="utf-8", errors="ignore") as f:
                            output = f.read()
                except Exception as e:
                    output = f"Error reading {path}: {e}"

        elif tool == "USE_BROWSER":
            try:
                from src.ai import run_browser_task
                import asyncio
                browser_task = response.get("input", response.get("task", response.get("tool_input", "")))
                output = asyncio.run(run_browser_task(task=browser_task, provider=provider, model_name=model))
            except ImportError as e:
                output = f"Error: Browser tool dependencies missing in this build. Details: {e}"
            except Exception as e:
                import traceback
                output = f"Error using browser: {str(e)}\n{traceback.format_exc()}"

        elif tool == "CLIPBOARD_MANAGER":
            import pyperclip
            action = response.get("action", response.get("input", response.get("tool_input", "")))
            if action == "read":
                try:
                    text = pyperclip.paste()
                    output = f"Clipboard content (read success): {text}"
                except Exception as e:
                    output = f"Error reading clipboard: {e}"
            elif action == "write":
                try:
                    content_to_write = response.get("content", response.get("text", response.get("data", "")))
                    pyperclip.copy(content_to_write)
                    verification = pyperclip.paste()
                    if verification == content_to_write:
                        output = f"Content successfully written and verified in clipboard: {content_to_write}"
                    else:
                        output = f"Warning: Written content did not match verification! Clipboard currently contains: {verification}"
                except Exception as e:
                    output = f"Error writing to clipboard: {e}"

        elif tool == "INTERPRET_SCREENSHOT":
            try:
                import base64
                from io import BytesIO
                from PIL import ImageGrab
                screenshot = ImageGrab.grab()
                buffered = BytesIO()
                screenshot.save(buffered, format="PNG")
                img_str = base64.b64encode(buffered.getvalue()).decode("utf-8")

                try:
                    with open(get_config_path("keys.json"), "r", encoding="utf-8") as f:
                        keys = json.load(f)
                except:
                    keys = {}

                output = get_image_description(img_str, provider, model, keys)
            except Exception as e:
                output = f"Error taking screenshot: {e}"

        elif tool == "GLOB":
            import glob
            pattern = response.get("pattern", response.get("input", response.get("tool_input", "")))
            if not pattern:
                output = "Error: No pattern provided."
            else:
                try:
                    matches = glob.glob(pattern, recursive=True)
                    if matches:
                        output = "Matches found:\n" + "\n".join(matches)
                    else:
                        output = f"No files matched the pattern: {pattern}"
                except Exception as e:
                    output = f"Error performing glob search: {e}"

        elif tool == "GREP":
            pattern = response.get("pattern", response.get("regex", response.get("query", response.get("input", response.get("tool_input", "")))))
            path = response.get("path", response.get("file", response.get("directory", "")))
            if not pattern or not path:
                output = "Error: Missing pattern or path."
            else:
                try:
                    results = []
                    if os.path.isfile(path):
                        files_to_check = [path]
                    elif os.path.isdir(path):
                        files_to_check = []
                        for root_dir, _, files in os.walk(path):
                            for f in files:
                                files_to_check.append(os.path.join(root_dir, f))
                    else:
                        files_to_check = []
                    
                    compiled = re.compile(pattern)
                    for fpath in files_to_check:
                        try:
                            with open(fpath, "r", encoding="utf-8") as f:
                                for i, line in enumerate(f, 1):
                                    if compiled.search(line):
                                        results.append(f"{fpath}:{i}:{line.strip()}")
                        except:
                            pass
                    
                    if results:
                        output = "Matches found:\n" + "\n".join(results[:100])
                        if len(results) > 100:
                            output += f"\n... and {len(results)-100} more matches."
                    else:
                        if not files_to_check:
                            output = f"Error: Path not found or invalid: {path}"
                        else:
                            output = "No matches found."
                except Exception as e:
                    output = f"Error performing grep search: {e}"

        else:
            output = f"Unknown tool: {tool}"

        if current_stop_event.is_set():
            break

        root.after(0, start_thinking_animation)
        response = llm.generate_response(user_input, validate_response=True, output=output,
                                         status_callback=update_status)

        # If we got a title now but previously had "New Chat" or no title, update it
        new_title = response.get("title")
        if not is_regenerate and new_title and current_chat_file:
            filename = os.path.basename(current_chat_file)
            if filename.startswith("New Chat") or filename.startswith("Chat_"):
                try:
                    new_title = new_title.strip().strip('"').strip("'")
                    safe_title = "".join([c for c in new_title if c.isalnum() or c in " _-"]).strip()
                    if safe_title:
                        new_path = os.path.join(get_chats_dir(), f"{safe_title}_{int(threading.current_thread().ident)}.json")
                        if current_chat_file != new_path:
                            if os.path.exists(current_chat_file):
                                os.rename(current_chat_file, new_path)
                            current_chat_file = new_path
                            root.after(0, update_sidebar_list)
                except Exception:
                    pass

        if current_stop_event.is_set():
            aborted = True
            break

    root.after(0, stop_thinking_animation)
    event = threading.Event()

    def _final_update():
        if my_stop_event is not current_stop_event:
            return  # Prevent older background threads from overwriting new UI
            
        if my_stop_event.is_set() or aborted:
            update_answer_box("[Task stopped by user]\n\n", "md_error", event)
            chat_history.append({"role": "ai", "content": "[Task stopped by user]"})
        else:
            update_answer_box(response.get("response", "") + "\n\n", None, event)
            chat_history.append({"role": "ai", "content": response.get("response", "")})
        
        save_current_chat()

    root.after(0, _final_update)
    event.wait()

    on_finish()

    if not current_stop_event.is_set():
        with open(get_config_path("memory.txt"), "a", encoding="utf-8") as mem_file:
            memory_content = response.get("memory")
            if memory_content:
                mem_file.write(f"{memory_content}\n")


def thread_handle_task():
    threading.Thread(target=handle_task, daemon=True).start()


def open_settings():
    from src.utils import get_config_path
    settings_win = ctk.CTkToplevel(root)
    settings_win.title("Settings")
    settings_win.geometry("450x480")
    settings_win.transient(root)

    current_settings = load_settings()
    enabled_tools_vars = {}

    all_tools = ["BASH", "POWERSHELL", "FILE_HANDLER", "EXEC_PY", "SEARCH_WEB", "READ_FILE", "USE_BROWSER", "CLIPBOARD_MANAGER",
                 "INTERPRET_SCREENSHOT", "GLOB", "GREP"]

    def on_settings_change(*args):
        active_tools = [t for t in all_tools if enabled_tools_vars[t].get()]
        save_settings(provider_var.get(), model_var.get(), mode_var.get(), active_tools)

    def on_model_change(choice):
        model_entry.delete(0, tk.END)
        model_entry.insert(0, choice)
        on_settings_change()

    def on_model_entry_change(event=None):
        model_var.set(model_entry.get())
        on_settings_change()


    provider_to_key_name = {
        "GitHubAI": "GITHUB_TOKEN",
        "Groq": "GROQ_API_KEY",
        "OpenRoute": "OPENROUTE_API_KEY",
        "Unclose": "UNCLOSE_API_KEY",
        "Anthropic": "ANTHROPIC_API_KEY",
        "OpenAI": "OPENAI_API_KEY",
        "Google": "GOOGLE_API_KEY",
        "Ollama": "OLLAMA_API_KEY",
        "LMStudio": "LMSTUDIO_API_KEY"
    }

    def load_keys():
        abs_keys_path = get_config_path("keys.json")
        try:
            with open(abs_keys_path, "r", encoding="utf-8") as f:
                return json.load(f)
        except (FileNotFoundError, json.JSONDecodeError):
            return {}

    def save_key(event=None):
        provider = provider_var.get()
        key_name = provider_to_key_name.get(provider)
        if key_name:
            keys = load_keys()
            keys[key_name] = api_key_entry.get()
            abs_keys_path = get_config_path("keys.json")
            with open(abs_keys_path, "w", encoding="utf-8") as f:
                json.dump(keys, f, indent=4)
            if event is None:
                update_status("Settings saved")
                settings_win.destroy()
        on_settings_change()

    def toggle_password():
        if api_key_entry.cget("show") == "*":
            api_key_entry.configure(show="")
            eye_button.configure(text="👁")
        else:
            api_key_entry.configure(show="*")
            eye_button.configure(text="🔒")

    def update_models(provider):
        models = {
            "GitHubAI": ["openai/gpt-4o", "microsoft/phi-4", "meta/llama-3.3-70b-instruct", "deepseek/deepseek-r1",
                         "openai/gpt-4o-mini", "meta/llama-4-maverick-17b-128e-instruct-fp8"],
            "Groq": ["openai/gpt-oss-120b", "llama-3.3-70b-versatile", "groq/compound", "qwen/qwen3-32b",
                     "meta-llama/llama-4-scout-17b-16e-instruct", "groq/compound-mini", "openai/gpt-oss-20b",
                     "llama-3.1-8b-instant", "openai/gpt-oss-safeguard-20b"],
            "OpenRoute": ["xiaomi/mimo-v2-pro", "anthropic/claude-4.6-sonnet", "minimax/minimax-m2.7",
                          "deepseek/deepseek-v3.2", "qwen/qwen3.6-plus:free", "anthropic/claude-4.6-opus",
                          "openai/gpt-5.4", "google/gemini-3.1-pro-preview", "moonshotai/kimi-k2.5",
                          "google/gemini-3.1-flash-lite-preview", "nousresearch/hermes-3-llama-3.1-405b:free",
                          "qwen/qwen3-coder:free", "google/gemma-3-27b-it:free", "openrouter/free",
                          "nvidia/nemotron-3-super-120b-a12b:free", "openai/gpt-oss-120b:free"],
            "Unclose": ["qwen3-vl:8b", "gpt-oss:latest", "deepseek-r1:14b-qwen-distill-q8_0"],
            "Anthropic": ["claude-opus-4-6", "claude-sonnet-4-6", "claude-opus-4-5", "claude-sonnet-4-5",
                          "claude-haiku-4-5", "claude-opus-4", "claude-sonnet-4", "claude-opus-4-5-20251101",
                          "claude-sonnet-4-5-20250929", "claude-haiku-4-5-20251001"],
            "OpenAI": ["gpt-5.4", "gpt-5.4-pro", "gpt-5.4-mini", "gpt-5.4-nano", "gpt-5-mini", "gpt-5", "gpt-5-nano",
                       "gpt-5.3-chat-latest", "gpt-4.1", "gpt-4o-mini"],
            "Google": ["gemini-3.1-pro-preview", "gemini-3-flash-preview", "gemini-3.1-flash-lite-preview",
                       "gemini-2.5-pro", "gemini-2.5-flash", "gemini-2.5-flash-lite", "gemma-4-31b-it",
                       "gemini-flash-latest"],
            "Ollama": ["llama3.1", "deepseek-r1", "llama3.2", "gemma3", "mistral", "qwen2.5", "qwen3", "llama3", "gemma2", "deepseek-r1:8b", "phi4"],
            "LMStudio": ["local-model", "qwen3-8b", "llama-3.1-8b", "mistral-7b-instruct", "gpt-oss-20b", "deepseek-r1-0528-qwen3-8b", "qwen3.5", "gemma4", "llama4", "mistral-small-3", "phi-4-mini", "lfm2-24b", "nemotron-3-super", "qwen3-coder"]
        }
        model_list = models.get(provider, ["default-model"])
        model_menu.configure(values=model_list)

        current_settings = load_settings()
        selected_model = ""
        if provider == current_settings.get("provider"):
            selected_model = current_settings.get("model")
        elif model_list:
            selected_model = model_list[0]
        
        if selected_model:
            model_var.set(selected_model)
            model_entry.delete(0, tk.END)
            model_entry.insert(0, selected_model)

        on_settings_change()

        keys = load_keys()
        key_name = provider_to_key_name.get(provider)
        api_key_entry.delete(0, tk.END)
        if key_name and key_name in keys:
            api_key_entry.insert(0, keys[key_name])

    # Left Column (Core Settings)
    main_frame = ctk.CTkFrame(settings_win, fg_color="transparent")
    main_frame.pack(fill="both", expand=True, padx=10, pady=10)

    left_col = ctk.CTkFrame(main_frame, fg_color="transparent")
    left_col.pack(side="left", fill="both", expand=True)

    ctk.CTkLabel(left_col, text="Provider:").pack(pady=(10, 0))
    provider_menu = ctk.CTkOptionMenu(left_col, variable=provider_var,
                                      values=["GitHubAI", "Groq", "OpenRoute", "Unclose", "Anthropic", "OpenAI",
                                              "Google", "Ollama", "LMStudio"], command=update_models)
    provider_menu.pack(pady=5)

    ctk.CTkLabel(left_col, text="Model:").pack(pady=(10, 0))
    model_menu = ctk.CTkOptionMenu(left_col, variable=model_var, values=[], command=on_model_change)
    model_menu.pack(pady=5)

    ctk.CTkLabel(left_col, text="Custom Model ID:").pack(pady=(5, 0))
    model_entry = ctk.CTkEntry(left_col, width=140)
    model_entry.pack(pady=5)
    model_entry.bind("<KeyRelease>", on_model_entry_change)

    ctk.CTkLabel(left_col, text="API Key:").pack(pady=(10, 0))
    key_frame = ctk.CTkFrame(left_col, fg_color="transparent")
    key_frame.pack(pady=5)

    api_key_entry = ctk.CTkEntry(key_frame, show="*", width=140)
    api_key_entry.pack(side="left")
    api_key_entry.bind("<KeyRelease>", save_key)

    eye_button = ctk.CTkButton(key_frame, text="🔒", width=30, command=toggle_password)
    eye_button.pack(side="left", padx=5)

    # Right Column (Tool Permissions)
    right_col = ctk.CTkFrame(main_frame, fg_color="#2a2a3a", corner_radius=10)
    right_col.pack(side="right", fill="both", expand=True, padx=(10, 0))

    ctk.CTkLabel(right_col, text="Tool Permissions", font=ctk.CTkFont(weight="bold")).pack(pady=(10, 5))

    scroll_frame = ctk.CTkScrollableFrame(right_col, fg_color="transparent", width=180, height=200)
    scroll_frame.pack(fill="both", expand=True, padx=5, pady=5)

    enabled_tools = current_settings.get("enabled_tools", all_tools)
    for tool_name in all_tools:
        var = tk.BooleanVar(value=(tool_name in enabled_tools))
        enabled_tools_vars[tool_name] = var
        cb = ctk.CTkCheckBox(scroll_frame, text=tool_name, variable=var, command=on_settings_change,
                             font=ctk.CTkFont(size=11))
        cb.pack(anchor="w", pady=2)

    save_button = ctk.CTkButton(settings_win, text="Close", command=settings_win.destroy)
    save_button.pack(pady=10)

    update_models(provider_var.get())


ctk.set_appearance_mode("dark")
ctk.set_default_color_theme("blue")

def _on_tkinter_error(exc, val, tb):
    import traceback
    err_str = str(val)
    if issubclass(exc, tk.TclError) and "invalid command name" in err_str and "!ctktextbox" in err_str:
        pass  # Ignore CustomTkinter textbox after loop error on destroy
    elif issubclass(exc, tk.TclError) and "invalid command name" in err_str and "!expandabletool" in err_str:
        pass
    else:
        traceback.print_exception(exc, val, tb)

root = ctk.CTk()
root.report_callback_exception = _on_tkinter_error
root.title("Prompt OS")
root.geometry("900x700")
root.minsize(500, 500)

from src.utils import resource_path
icon_path = resource_path("assets", "logo.ico")
if os.path.exists(icon_path):
    try:
        from PIL import Image, ImageTk
        root.iconbitmap(icon_path)
        img = ImageTk.PhotoImage(Image.open(icon_path))
        root.wm_iconphoto(True, img)
    except Exception:
        pass

initial_settings = load_settings()
provider_var = ctk.StringVar(value=initial_settings.get("provider", "GitHubAI"))
model_var = ctk.StringVar(value=initial_settings.get("model", "openai/gpt-4o-mini"))
mode_var = ctk.StringVar(value=initial_settings.get("mode", "FAST"))

root.configure(padx=0, pady=0)
root.configure(fg_color="#0f0f1a")

main_container = ctk.CTkFrame(root, fg_color="transparent")
main_container.pack(fill="both", expand=True)

sidebar_frame = ctk.CTkFrame(main_container, width=200, corner_radius=0, fg_color="#1a1a24")
sidebar_frame.pack(side="left", fill="y")
sidebar_frame.pack_propagate(False)

content_frame = ctk.CTkFrame(main_container, fg_color="#0f0f1a", corner_radius=0)
content_frame.pack(side="right", fill="both", expand=True)

from PIL import Image
logo_pil = Image.open(resource_path("assets", "logo.png"))
logo_img = ctk.CTkImage(light_image=logo_pil, dark_image=logo_pil, size=(32, 32))
logo_label = ctk.CTkLabel(sidebar_frame, text=" Prompt OS", image=logo_img, compound="left", font=ctk.CTkFont(size=16, weight="bold"))
logo_label.pack(side="top", anchor="w", padx=15, pady=15)

new_chat_btn = ctk.CTkButton(sidebar_frame, text="+ New Chat", height=32, font=ctk.CTkFont(weight="bold"), 
                             fg_color="#2a2a3a", hover_color="#3a3a5a", anchor="w", command=clear_chat)
new_chat_btn.pack(side="top", pady=(10, 5), padx=15, fill="x")

history_scroll = ctk.CTkScrollableFrame(sidebar_frame, fg_color="transparent", label_text="Recent Chats", 
                                         label_text_color="#666666", label_font=ctk.CTkFont(size=11, weight="bold"))
history_scroll.pack(fill="both", expand=True, padx=5, pady=5)

settings_btn = ctk.CTkButton(sidebar_frame, text="⚙ Settings", height=32, font=ctk.CTkFont(weight="bold"), 
                             fg_color="transparent", hover_color="#2a2a3a", anchor="w", command=open_settings)
settings_btn.pack(side="bottom", pady=15, padx=15, fill="x")

mode_selection = ctk.CTkSegmentedButton(sidebar_frame, values=["FAST", "THINKING", "PRO"], variable=mode_var,
                                        command=lambda m: save_settings(provider_var.get(), model_var.get(), m))
mode_selection.pack(side="bottom", padx=15, pady=(0, 10), fill="x")

header_frame = ctk.CTkFrame(content_frame, height=44, fg_color="#1a1a24", corner_radius=0)
header_frame.pack(fill="x", side="top")

model_indicator_var = ctk.StringVar(value=f"{model_var.get()} · {mode_var.get()}")
def update_model_indicator(*args):
    model_indicator_var.set(f"{model_var.get()} · {mode_var.get()}")
model_var.trace_add("write", update_model_indicator)
mode_var.trace_add("write", update_model_indicator)

model_chip = ctk.CTkLabel(header_frame, textvariable=model_indicator_var, font=ctk.CTkFont(size=12, weight="bold"),
                          fg_color="#2a2a3a", text_color="#a9a9c5", corner_radius=6, padx=10)
model_chip.pack(side="right", padx=20, pady=8)
content_frame.pack(fill="both", expand=True)

ai_answer_box = ctk.CTkTextbox(
    content_frame, font=ctk.CTkFont(size=14, family="Helvetica"),
    corner_radius=12,
    fg_color="#1e1e2e",
    border_color="#3a3a55",
    border_width=1,
    text_color="#e0e0e0"
)
ai_answer_box.pack(pady=(15, 10), padx=20, fill="both", expand=True)

toolbar_frame = ctk.CTkFrame(content_frame, fg_color="transparent", height=24)
toolbar_frame.pack(fill="x", padx=20, pady=(0, 5))

clear_btn = ctk.CTkButton(toolbar_frame, text="Clear Chat", width=60, height=24, fg_color="transparent", text_color="#a9a9c5", hover_color="#2a2a3a", command=clear_chat)
clear_btn.pack(side="left")

status_label = ctk.CTkLabel(
    toolbar_frame,
    text="Ready",
    font=ctk.CTkFont(size=12, slant="italic"),
    text_color="#888888"
)
status_label.pack(side="left", padx=10)

regen_btn = ctk.CTkButton(toolbar_frame, text="Regenerate", width=60, height=24, fg_color="transparent", text_color="#a9a9c5", hover_color="#2a2a3a", command=regenerate)
regen_btn.pack(side="right")

copy_btn = ctk.CTkButton(toolbar_frame, text="Copy Output", width=60, height=24, fg_color="transparent", text_color="#a9a9c5", hover_color="#2a2a3a", command=copy_last)
copy_btn.pack(side="right", padx=(0, 5))

input_frame = ctk.CTkFrame(content_frame, fg_color="transparent")
input_frame.pack(fill="x", padx=20, pady=(0, 15))

text_input = ctk.CTkTextbox(input_frame, font=ctk.CTkFont(size=14), height=40, wrap="word", border_width=2, border_color="#5a5a75", corner_radius=8, fg_color="#1a1a26")
text_input.pack(side="left", fill="x", expand=True)

send_button = ctk.CTkButton(input_frame, text="Send", width=80, height=40, font=ctk.CTkFont(weight="bold"), fg_color="#2a2a3a", text_color="#a9a9c5", hover_color="#3a3a5a", command=thread_handle_task)
send_button.pack(side="right", padx=(10, 0), anchor="s")

stop_button = ctk.CTkButton(input_frame, text="Stop", width=80, height=40, fg_color="#c42b2b", hover_color="#a82525", command=stop_task)

def update_send_btn_state():
    if text_input.get("1.0", "end-1c").strip():
        send_button.configure(fg_color="#1f538d", text_color="#ffffff", hover_color="#14375e")
    else:
        send_button.configure(fg_color="#2a2a3a", text_color="#a9a9c5", hover_color="#3a3a5a")

def resize_textbox_event(event=None):
    resize_textbox()
    update_send_btn_state()

text_input.bind("<KeyRelease>", resize_textbox_event)

def on_enter(event):
    if not event.state & 0x0001:
        thread_handle_task()
        return "break"

text_input.bind("<Return>", on_enter)

def on_ctrl_k(event):
    clear_chat()
    return "break"
root.bind("<Control-k>", on_ctrl_k)

configure_markdown_tags(ai_answer_box)
ai_answer_box.configure(state="disabled")

# Delay loading the sidebar to make the main window appear faster
root.after(100, update_sidebar_list)

root.mainloop()
